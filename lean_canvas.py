# -*- coding: utf-8 -*-
"""
精益画布 PPT 生成器（V2 重写版）
基于大赛统一模板填充 risk_agent_v2 项目内容，与实际代码严格对应。
模板占位符（Slide 1）：
  id=11 问题描述      id=12 解决方案      id=13 研究方案
  id=15 技术亮点      id=16 落地可行性    id=17 缺点及改进
  id=18 揭榜编号      id=3  团队名        id=14 指导老师
Slide 0 标题页：
  id=8  作品标题      id=18 揭榜编号
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

TEMPLATE_PATH = r"D:/08-智能风控与量化建模赛道-江苏银行-信贷风控模型智能监测与自主迭代/附件3：第五届中国研究生金融科技创新大赛精益画布模板.pptx"
OUTPUT_PATH = r"D:/Work Buddy/risk_agent/risk_agent_v2/精益画布_v2.pptx"

prs = Presentation(TEMPLATE_PATH)

# ============ Slide 0: 标题页 ============
slide0 = prs.slides[0]
for shape in slide0.shapes:
    if shape.shape_id == 8:
        shape.text_frame.text = "信贷风控模型智能监测与自主迭代系统"
    elif shape.shape_id == 18:
        shape.text_frame.text = "揭榜挂帅-14"

# ============ Slide 1: 精益画布内容页 ============
slide1 = prs.slides[1]

content_map = {
    # 11 — 问题描述
    11: """面向信贷风控模型全生命周期的三大技术卡点：

1. 性能衰减难察觉：传统监测仅看 KS/AUC 宏观指标，无法区分特征漂移、标签漂移、联合漂移等不同根因，迭代策略"一刀切"，轻漂移过度重训、重漂移延误修复。

2. 标签延迟 30 天的决策盲区：贷前 A 卡坏账标签需等还款表现期，数据不全时无法判断模型是否健康，导致修复决策延误 1 个月，坏账损失扩大。

3. 迭代决策依赖人工经验：何时迭代、选什么策略、上线后如何验证，全靠风控分析师拍脑袋，缺乏可审计、可回放的标准化决策闭环。""",

    # 12 — 解决方案
    12: """核心思路：构建"监测→根因→迭代→部署"五 Agent 自主闭环，替代人工盯盘。

M1 数据与漂移注入 → M2 双阶段监测 → M3 双层根因+LLM 四维归因 → M4 三策略迭代 → M5 四层灰度部署 → Champion/Challenger 动态替换。

一键入口 run_demo.py，五场景并行验证（1 次真实 OOT + 4 个平行月度场景 A/B/C/D），全程 JSON+DOCX 留痕。""",

    # 13 — 研究方案（按代码实现严格对应）
    13: """研究方案（与代码模块一一对应）：

1. M1 数据层：四步预处理——缺失校验→异常值截断→共线性剔除→强特征剔除，34→31 特征。基线模型 XGB/LGB/LR 三家选优，最终 LightGBM 胜出。时序场景生成：从 train 2025-12 采样每月 2 万条，按 SCENARIO_SPEC 注入特征漂移与标签翻转。

2. M2 监测层：双阶段监测——标签未成熟期（无标签 4 项规则）+ 标签成熟期（有标签 R1~R7 + DQ1~DQ2）。LOW/MEDIUM/HIGH 三级告警，取最高级为综合告警等级。

3. M3 诊断层：双层根因——L1（5 项指标级诊断）+ L2（5 项特征级归因）。DeepSeek-V4-Flash LLM 四维归因，证据隔离设计，结构化输出含独立异常计数、首要根因、策略推荐、证据引用。

4. M4 迭代层：三策略——light（增量微调）/ standard（特征剔除 + 先验重加权）/ major（重度重构 + 延迟重训），对应不同漂移强度分级响应。

5. M5 部署层：四层灰度——L1 离线全量回放（KS 恢复率≥0.90 + AUC 非劣 0.015 + PSI<0.30 + 合格线）→ L2-20%（护栏：KS 比率≥0.95 + AUC 非劣 + PSI<0.30 + 预测坏账率≤1.2×champion）→ L2-50%（加业务指标：通过率偏差<5pp + 坏账捕获率 Top10% 不显著降）→ 全量（DeLong 显著性检验 + KS Bootstrap 500 次，显著优/非劣/显著劣三判定）。防泄漏：major 策略 relabel_ids 从评测集自动剔除。显著优或非劣即替换 champion 并写 champion_registry.json。""",

    # 15 — 技术亮点
    15: """技术亮点（严格对齐代码实现）：

1. 双阶段监测：独创"标签未成熟期无标签监测 + 成熟期有标签监测"双轨制，解决标签延迟 30 天的决策盲区，适应窗口不误报、不空转。

2. 噪声底机制：8 次×20000 条重采样取 95% 分位数作为 IV 衰减/重要性偏移的噪声阈值，过滤采样噪声，避免小样本波动误触发迭代。

3. LLM 自主推理：大模型融合根因知识库，自主推理四维归因与迭代策略推荐，证据隔离设计防锚定偏差，规则引擎与 AI 双通道交叉验证。

4. 防泄漏机制：major 策略延迟重训时 relabel_ids 落盘，M5 评测时自动剔除已参与训练的行，杜绝自评作弊。

5. 三策略自适应：light/standard/major 对应不同漂移强度，训练成本从分钟级到小时级分级，避免过度重训。选优门控拆分过拟合 gap 与漂移适应度，不混淆两类问题。

6. 灰度发布四阶段 + 显著性检验：L1/L2-20%/L2-50%/全量，每阶段独立护栏任一不过即回滚；全量阶段 DeLong 检验 + KS Bootstrap 500 次，显著优替换、非劣替换、显著劣 hold，三态可审计。""",

    # 16 — 落地可行性
    16: """落地可行性：

1. 技术可行性：全栈 Python 开源生态（XGBoost/LightGBM/scikit-learn/python-docx），无自研框架依赖，银行现有技术栈可直接复用；模型可解释性强（树模型 + LR），满足银保监会模型风险管理要求。

2. 数据可行性：仅需贷前申请数据（id_card/apply_time/is_bad + 31 特征），不依赖外部征信数据，银行内部数据即可跑通；34 特征零缺失，数据质量要求低。

3. 部署可行性：一键入口 run_demo.py 支持 --all/--scenario/--oot/--prepare 四种模式，标准环境 pip install -r requirements.txt 即可运行，无需 GPU；输出全部 JSON+DOCX 留痕，可审计可回放。

4. 成本可行性：三策略分级设计，轻度漂移用 light（分钟级），重度才用 major（小时级），算力成本可控；LLM 调用可选 DeepSeek API 或降级为规则引擎兜底。

5. 验证充分性：1 次真实 OOT（2026-01 test 2 万条）+ 4 个平行月度场景（A 无漂移/B 轻量特征漂移/C 中度标签漂移/D 重度联合漂移），共 5 轮×3 月=15 次决策，9 次成功上线（4 显著优替换 + 5 非劣替换），0 回滚。""",

    # 17 — 缺点及改进方案
    17: """缺点及改进方案：

1. 分群维度有限：当前仅 4 个分群维度（age/city_tier/loan_amount/repayment_period）。改进：接入渠道/产品/地区等业务维度，支持自定义分群与分群级漂移监测。

2. LLM 依赖外部 API：DeepSeek 调用存在网络延迟与 API 稳定性风险，已有规则引擎兜底但 AI 归因质量受限。改进：支持本地部署开源模型（Qwen/ChatGLM），或用蒸馏小模型替代。

3. 时序场景为合成数据：4 个平行场景的漂移为从 train 采样后注入，非真实线上数据。改进：接入真实线上流量回放，验证漂移注入与真实漂移的分布一致性。

4. 迭代策略为预设规则：light/standard/major 三策略及触发阈值依赖人工标定，尚未实现策略选择的自动优化。改进：引入历史迭代效果反馈，让系统自动学习最优策略组合与阈值参数。

5. 监测指标以全局聚合为主：当前 KS/AUC/PSI 均为全量客群聚合值，细分客群的局部漂移可能被均值稀释。改进：增加分群级指标监测，对高价值客群设置更敏感的独立告警阈值。""",
}

for shape in slide1.shapes:
    if shape.shape_id in content_map:
        text = content_map[shape.shape_id]
        shape.text_frame.text = text
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(9)
                run.font.name = "宋体"

# Slide 1 团队名 / 指导老师 / 揭榜编号
for shape in slide1.shapes:
    if shape.shape_id == 3:
        shape.text_frame.text = "揭榜挂帅-14 团队"
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(10)
                r.font.name = "宋体"
    elif shape.shape_id == 18:
        shape.text_frame.text = "揭榜挂帅-14"
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(10)
                r.font.name = "宋体"
    elif shape.shape_id == 14:
        shape.text_frame.text = "指导老师"
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(10)
                r.font.name = "宋体"

prs.save(OUTPUT_PATH)
print(f"Saved: {OUTPUT_PATH}")
